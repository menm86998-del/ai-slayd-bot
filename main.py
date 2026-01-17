import telebot
from pptx import Presentation
from pptx.util import Inches, Pt
import requests
import os
import io
from g4f.client import Client # Bepul AI kutubxonasi

# 1. SOZLAMALAR
BOT_TOKEN = "8462276425:AAEgRcsPSia3PxqTcyg2rn99ef92mfHg0fs"
bot = telebot.TeleBot(BOT_TOKEN)
client = Client()

# 2. BEPUL AI ORQALI MATN YARATISH
def get_ai_content(topic):
    try:
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": f"Menga '{topic}' mavzusida 15 ta slayd uchun reja ber. Har bir slayd sarlavhasi va 3 qator matni bo'lsin. O'zbek tilida."}]
        )
        return response.choices[0].message.content
    except:
        return f"{topic} haqida ma'lumot\nBu yerga matn kiritiladi."

# 3. SLAYDNI YIG'ISH
def create_presentation(topic, filename):
    # O'zgarish shu yerda: Bo'sh prezentatsiya emas, shablonni ochamiz
    if os.path.exists("template.pptx"):
        prs = Presentation("template.pptx")
    else:
        prs = Presentation() # Agar shablon bo'lmasa, oddiy ochiq qoladi
    
    # AI matnni oladi
    content = get_ai_content(topic)
    # Slaydlarni qismlarga bo'lish
    sections = content.split('\n\n')

    for i in range(15):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Sarlavha va matnni kiritish
        if i < len(sections):
            text_parts = sections[i].split('\n')
            slide.shapes.title.text = text_parts[0]
            slide.placeholders[1].text = "\n".join(text_parts[1:])
        else:
            slide.shapes.title.text = f"{topic} - Davomi"
            slide.placeholders[1].text = "Qo'shimcha ma'lumotlar..."

        # Har bir sahifaga mavzuga mos rasm qo'yish
        try:
            img_url = f"https://source.unsplash.com/featured/800x600?{topic.replace(' ', ',')},{i}"
            img_data = requests.get(img_url, timeout=10).content
            img_stream = io.BytesIO(img_data)
            slide.shapes.add_picture(img_stream, Inches(5.5), Inches(2), width=Inches(4))
        except:
            pass

    prs.save(filename)

@bot.message_handler(func=lambda message: True)
def handle_message(message):
    topic = message.text
    msg = bot.reply_to(message, "🚀 Gamma AI (Bepul versiya) ishga tushdi. 15 betlik slayd tayyorlanmoqda, kuting...")
    
    file_name = f"{topic}.pptx"
    try:
        create_presentation(topic, file_name)
        with open(file_name, 'rb') as doc:
            bot.send_document(message.chat.id, doc, caption=f"✅ '{topic}' mavzusidagi 15 betlik slayd tayyor!")
        os.remove(file_name)
    except Exception as e:
        bot.reply_to(message, f"Xatolik: {str(e)}")

print("Bot muvaffaqiyatli ishga tushdi (Limitsiz versiya)...")
bot.polling()






